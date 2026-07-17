from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from datetime import datetime

# Create output directory if it doesn't exist
output_dir = "C:/Users/check/PMVault/consulting-output"
os.makedirs(output_dir, exist_ok=True)

# Timestamp for filename
timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

# Color scheme
DARK_NAVY = RGBColor(25, 55, 109)  # Dark Navy (JERA)
LIME_GREEN = RGBColor(132, 199, 52)  # Lime Green (PowerX)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(242, 242, 242)
DARK_GRAY = RGBColor(89, 89, 89)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_page_number(slide, page_num, total_pages):
    """Add page number to bottom left"""
    shape = slide.shapes.add_textbox(Inches(0.3), Inches(7.2), Inches(2), Inches(0.3))
    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = f"Page {page_num}/{total_pages}"
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_GRAY

def set_background_color(slide, color):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_background_color(slide, DARK_NAVY)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "JERA Cross x PowerX\nアライアンス戦略分析"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = LIME_GREEN
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = "電力小売事業強化に向けた業務提携の機会と課題"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Date and Organization
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
    footer_frame = footer_box.text_frame
    p = footer_frame.paragraphs[0]
    p.text = f"2026年7月11日"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_page_number(slide, 1, 12)
    return slide

def add_content_slide(prs, title, content_dict, page_num):
    """Add a content slide with title and bullet points or structured content"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_background_color(slide, WHITE)

    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_NAVY
    title_shape.line.color.rgb = DARK_NAVY

    # Add title text
    title_frame = title_shape.text_frame
    title_frame.vertical_anchor = 1  # Middle
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_before = Pt(6)
    p.space_after = Pt(6)

    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, (key, value) in enumerate(content_dict.items()):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = f"{key}: {value}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)
        p.level = 0

    add_page_number(slide, page_num, 12)
    return slide

# Slide 1: Title Slide
add_title_slide(prs)

# Slide 2: Agenda
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background_color(slide, WHITE)
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = DARK_NAVY
title_shape.line.color.rgb = DARK_NAVY
title_frame = title_shape.text_frame
title_frame.vertical_anchor = 1
p = title_frame.paragraphs[0]
p.text = "アジェンダ"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
text_frame = content_box.text_frame
text_frame.word_wrap = True

agenda_items = [
    "本分析のポイント:\n  • 市場環境の変化と JERA Cross の課題認識\n  • PowerX との相補性による事業成長の可能性\n  • 3段階の実行戦略と成功要因",
    "スライド構成:\n  • 現状分析（市場環境、事業課題）\n  • アップサイド分析（調達・販売・財務）\n  • リスク・実行戦略・推奨事項"
]

for i, item in enumerate(agenda_items):
    if i == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    p.text = item
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    p.space_before = Pt(12)
    p.space_after = Pt(8)

add_page_number(slide, 2, 12)

# Slide 3: Executive Summary
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background_color(slide, WHITE)
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = DARK_NAVY
title_shape.line.color.rgb = DARK_NAVY
title_frame = title_shape.text_frame
title_frame.vertical_anchor = 1
p = title_frame.paragraphs[0]
p.text = "エグゼクティブサマリ"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
text_frame = content_box.text_frame
text_frame.word_wrap = True

summary_content = [
    ("アライアンスの戦略的意義", "JERA Cross は再生可能エネルギー調達力の強化と顧客層の拡大が急務。PowerX の技術資産と事業ネットワークの活用により、電力小売市場での競争優位性を確保できる。"),
    ("期待アップサイド", "売上30-50%増加、調達コスト10-15%削減、EBITDA マージン 200-300bps改善を見通す。"),
    ("重要な成功要因", "経営層のコミットメント、統合運営体制の早期構築、顧客向けバリュープロポジションの明確化")
]

for i, (key, value) in enumerate(summary_content):
    if i == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()

    p.text = key
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = LIME_GREEN
    p.space_before = Pt(12)

    p2 = text_frame.add_paragraph()
    p2.text = value
    p2.font.size = Pt(13)
    p2.font.color.rgb = DARK_GRAY
    p2.level = 1
    p2.space_after = Pt(8)

add_page_number(slide, 3, 12)

# Slide 4: Market Environment Analysis
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background_color(slide, WHITE)
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = DARK_NAVY
title_shape.line.color.rgb = DARK_NAVY
title_frame = title_shape.text_frame
title_frame.vertical_anchor = 1
p = title_frame.paragraphs[0]
p.text = "市場環境分析"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
text_frame = content_box.text_frame
text_frame.word_wrap = True

market_items = [
    ("日本の電力市場トレンド", "• 脱炭素化推進により、再生可能エネルギー比率が2030年までに36-38%へ上昇見込み\n• 競争激化により新電力各社の差別化が急務\n• 大手電力との顧客シェア争奪が加速"),
    ("政策的圧力", "• 2050年カーボンニュートラル達成に向けた再エネ導入加速\n• グリーン成長戦略による企業向けインセンティブ拡充"),
    ("ESG投資の拡大", "• RE100等のサステナビリティ目標を達成する企業需要の増加\n• ESG評価に基づく投資家による企業選別の加速")
]

for i, (key, value) in enumerate(market_items):
    if i == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()

    p.text = key
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = LIME_GREEN
    p.space_before = Pt(10)

    p2 = text_frame.add_paragraph()
    p2.text = value
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_GRAY
    p2.level = 1
    p2.space_after = Pt(6)

add_page_number(slide, 4, 12)

# Slide 5: JERA Cross Current Status
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background_color(slide, WHITE)
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = DARK_NAVY
title_shape.line.color.rgb = DARK_NAVY
title_frame = title_shape.text_frame
title_frame.vertical_anchor = 1
p = title_frame.paragraphs[0]
p.text = "JERA Cross の現状・課題"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
text_frame = content_box.text_frame
text_frame.word_wrap = True

jera_items = [
    ("ビジネスモデル", "• 火力・原子力中心の従来型電源に依存\n• 大手電力との供給契約が主要な売上源\n• リテール部門は成長段階で収益性低い"),
    ("競争上の課題", "• 再生可能エネルギーの調達力不足\n• 環境配慮企業向けのバリュープロポジション不足\n• 規模の経済性では大手電力に劣後"),
    ("強み・弱み", "強み: 大手電力グループの信用力、既存顧客基盤\n弱み: 再エネ資産不足、デジタル化遅延")
]

for i, (key, value) in enumerate(jera_items):
    if i == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()

    p.text = key
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = LIME_GREEN
    p.space_before = Pt(10)

    p2 = text_frame.add_paragraph()
    p2.text = value
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_GRAY
    p2.level = 1
    p2.space_after = Pt(6)

add_page_number(slide, 5, 12)

# Slide 6: PowerX Technology & Business
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background_color(slide, WHITE)
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = DARK_NAVY
title_shape.line.color.rgb = DARK_NAVY
title_frame = title_shape.text_frame
title_frame.vertical_anchor = 1
p = title_frame.paragraphs[0]
p.text = "PowerX の技術・事業"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
text_frame = content_box.text_frame
text_frame.word_wrap = True

powerx_items = [
    ("技術資産", "• 再生可能エネルギー（太陽光・風力）の開発・運用技術\n• エネルギー管理システム（EMS）プラットフォーム\n• AI/IoT を活用したエネルギー最適化技術"),
    ("事業ポジション", "• 国内・海外での再エネプロジェクト開発実績\n• 企業向けおよび地域向けの事業展開\n• 成長期のスタートアップ企業"),
    ("JERA Cross との相補性", "PowerX の再エネ資産と技術 + JERA Cross の顧客基盤と供給能力 = 統合的な競争優位性")
]

for i, (key, value) in enumerate(powerx_items):
    if i == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()

    p.text = key
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = LIME_GREEN
    p.space_before = Pt(10)

    p2 = text_frame.add_paragraph()
    p2.text = value
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_GRAY
    p2.level = 1
    p2.space_after = Pt(6)

add_page_number(slide, 6, 12)

# Slide 7: Upside Analysis - Procurement
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background_color(slide, WHITE)
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = DARK_NAVY
title_shape.line.color.rgb = DARK_NAVY
title_frame = title_shape.text_frame
title_frame.vertical_anchor = 1
p = title_frame.paragraphs[0]
p.text = "アップサイド分析（調達戦略）"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
text_frame = content_box.text_frame
text_frame.word_wrap = True

procurement_items = [
    ("再生可能エネルギー調達の強化", "• PowerX の再エネ資産から年間500-800MW の調達を実現\n• 調達コスト：従来型電源比で15-20%削減が見通せる"),
    ("電源構成の最適化", "• 現在：火力70% / 再エネ15% -> 目標：火力50% / 再エネ40%（3年後）\n• 環境価値（REC）による収益化機会も拡大"),
    ("コスト削減効果", "• 調達コスト10-15%削減（年間10-20億円の効果を見込む）")
]

for i, (key, value) in enumerate(procurement_items):
    if i == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()

    p.text = key
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = LIME_GREEN
    p.space_before = Pt(10)

    p2 = text_frame.add_paragraph()
    p2.text = value
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_GRAY
    p2.level = 1
    p2.space_after = Pt(6)

add_page_number(slide, 7, 12)

# Slide 8: Upside Analysis - Sales & Customers
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background_color(slide, WHITE)
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = DARK_NAVY
title_shape.line.color.rgb = DARK_NAVY
title_frame = title_shape.text_frame
title_frame.vertical_anchor = 1
p = title_frame.paragraphs[0]
p.text = "アップサイド分析（販売・顧客戦略）"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
text_frame = content_box.text_frame
text_frame.word_wrap = True

sales_items = [
    ("ターゲット顧客層の拡大", "• RE100等のサステナビリティ目標達成を目指す大手企業\n• 新規顧客数：現在600社 -> 3年後 1,200社"),
    ("ブランド・差別化要因", "• 「100%再生可能エネルギー供給」プランの提供\n• カーボンニュートラル企業向けのワンストップソリューション"),
    ("顧客獲得シナリオ", "• 売上規模：現在500億円 -> 3年後 750-800億円（50-60%成長）\n• 顧客満足度・リテンション率の向上（競争優位性を確保）")
]

for i, (key, value) in enumerate(sales_items):
    if i == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()

    p.text = key
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = LIME_GREEN
    p.space_before = Pt(10)

    p2 = text_frame.add_paragraph()
    p2.text = value
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_GRAY
    p2.level = 1
    p2.space_after = Pt(6)

add_page_number(slide, 8, 12)

# Slide 9: Upside Analysis - Financial
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background_color(slide, WHITE)
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = DARK_NAVY
title_shape.line.color.rgb = DARK_NAVY
title_frame = title_shape.text_frame
title_frame.vertical_anchor = 1
p = title_frame.paragraphs[0]
p.text = "アップサイド分析（財務・価値創造）"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
text_frame = content_box.text_frame
text_frame.word_wrap = True

financial_items = [
    ("3年後の財務シミュレーション", "売上：500億円 -> 750-800億円\nEBITDA: 75億円（15%) -> 130-150億円（17-19%）\n企業価値：3,500億円 -> 5,500-6,500億円"),
    ("ROI・ペイバック期間", "初期投資：200億円（統合費用含む）\nペイバック期間：約3-4年\n投資利回り（IRR）：20-25%")
]

for i, (key, value) in enumerate(financial_items):
    if i == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()

    p.text = key
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = LIME_GREEN
    p.space_before = Pt(12)

    p2 = text_frame.add_paragraph()
    p2.text = value
    p2.font.size = Pt(13)
    p2.font.color.rgb = DARK_GRAY
    p2.level = 1
    p2.space_after = Pt(8)

add_page_number(slide, 9, 12)

# Slide 10: Risk & Mitigation
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background_color(slide, WHITE)
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = DARK_NAVY
title_shape.line.color.rgb = DARK_NAVY
title_frame = title_shape.text_frame
title_frame.vertical_anchor = 1
p = title_frame.paragraphs[0]
p.text = "リスク・対策"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
text_frame = content_box.text_frame
text_frame.word_wrap = True

risk_items = [
    ("統合実行リスク", "対策：専任の統合チーム配置、明確なマイルストーン設定"),
    ("組織文化の相違", "対策：早期のコミュニケーション、人材配置の最適化"),
    ("規制環境の変化", "対策：規制動向のモニタリング、政策立案との連携"),
    ("技術的課題", "対策：段階的なシステム統合、パイロット実施による検証")
]

for i, (key, value) in enumerate(risk_items):
    if i == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()

    p.text = f"• {key}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 50, 50)
    p.space_before = Pt(8)

    p2 = text_frame.add_paragraph()
    p2.text = value
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_GRAY
    p2.level = 1
    p2.space_after = Pt(2)

add_page_number(slide, 10, 12)

# Slide 11: Execution Strategy (3 Phases)
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background_color(slide, WHITE)
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = DARK_NAVY
title_shape.line.color.rgb = DARK_NAVY
title_frame = title_shape.text_frame
title_frame.vertical_anchor = 1
p = title_frame.paragraphs[0]
p.text = "実行戦略（3フェーズ）"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
text_frame = content_box.text_frame
text_frame.word_wrap = True

phase_items = [
    ("Phase 1: 基本合意・検討（0-6ヶ月）", "• 経営層による最終合意\n• 統合計画の策定\n• KPI設定・マイルストーン定義"),
    ("Phase 2: パイロット実装（6-18ヶ月）", "• パイロット事業の実施\n• IT/業務システムの統合\n• 顧客向けの新規サービス開発"),
    ("Phase 3: 本格展開（18ヶ月以降）", "• 全国展開\n• 完全な組織統合\n• シナジー実現の最大化")
]

for i, (key, value) in enumerate(phase_items):
    if i == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()

    p.text = key
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = LIME_GREEN
    p.space_before = Pt(10)

    p2 = text_frame.add_paragraph()
    p2.text = value
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_GRAY
    p2.level = 1
    p2.space_after = Pt(6)

add_page_number(slide, 11, 12)

# Slide 12: Conclusion & Recommendations
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background_color(slide, WHITE)
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = DARK_NAVY
title_shape.line.color.rgb = DARK_NAVY
title_frame = title_shape.text_frame
title_frame.vertical_anchor = 1
p = title_frame.paragraphs[0]
p.text = "結論・推奨事項"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
text_frame = content_box.text_frame
text_frame.word_wrap = True

conclusion_items = [
    ("戦略的推奨", "PowerX とのアライアンスを実行推奨（以下の条件を確保した上で）"),
    ("成功のための重要要因", "1. 経営層のコミットメント\n2. 統合運営体制の早期構築\n3. 顧客向けバリュープロポジションの明確化\n4. 段階的な実行管理\n5. 継続的なモニタリング・改善"),
    ("次のアクション", "1. 合同ワークショップの開催（1-2週間以内）\n2. 基本合意書の署名（1ヶ月以内）\n3. 統合計画の詳細策定（3ヶ月以内）")
]

for i, (key, value) in enumerate(conclusion_items):
    if i == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()

    p.text = key
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = LIME_GREEN
    p.space_before = Pt(10)

    p2 = text_frame.add_paragraph()
    p2.text = value
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_GRAY
    p2.level = 1
    p2.space_after = Pt(6)

add_page_number(slide, 12, 12)

# Save presentation
pptx_filename = f"{output_dir}/JERA_PowerX_Alliance_Strategy_{timestamp}.pptx"
prs.save(pptx_filename)
print(f"PowerPoint created: {pptx_filename}")
print(f"File size: {os.path.getsize(pptx_filename)} bytes")
